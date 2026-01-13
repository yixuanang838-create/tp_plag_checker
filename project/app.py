import streamlit as st
import tempfile
import zipfile
import itertools
import re
import docx2txt
import pptx
from pdfminer.high_level import extract_text as extract_pdf_text
from sentence_transformers import SentenceTransformer, util
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
import pandas as pd
from io import BytesIO
import os

# =====================================================
# PAGE CONFIG
# =====================================================
st.set_page_config(
    page_title="Plagiarism Checker",
    layout="centered"
)

# =====================================================
# HEADER
# =====================================================
col1, col2 = st.columns([1, 6])
with col1:
    image_path = os.path.join(os.path.dirname(__file__), "sonic-bear.png")
st.image(image_path, width=90)
with col2:
    st.markdown(
        """
        <div style="line-height:1.2;">
            <h3 style="margin-bottom:4px;">
                School of Informatics and IT
            </h3>
            <p style="margin-top:0; font-size:16px; color:#444;">
                Academic Integrity – Plagiarism Detection System
            </p>
        </div>
        """,
        unsafe_allow_html=True
    )
st.markdown("---")

# =====================================================
# LOAD MODEL
# =====================================================
@st.cache_resource(show_spinner=False)
def load_model():
    return SentenceTransformer("all-MiniLM-L6-v2")

model = load_model()

# =====================================================
# TEXT EXTRACTION (CACHED)
# =====================================================
@st.cache_data(show_spinner=False)
def extract_file_text(file_bytes, filename):
    from io import BytesIO
    f = BytesIO(file_bytes)
    f.name = filename.lower()
    name = f.name
    if name.endswith(".txt"):
        return f.read().decode("utf-8", errors="ignore")
    if name.endswith(".pdf"):
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(f.read())
            tmp.flush()
            return extract_pdf_text(tmp.name)
    if name.endswith(".docx"):
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(f.read())
            tmp.flush()
            return docx2txt.process(tmp.name)
    if name.endswith(".pptx"):
        prs = pptx.Presentation(f)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    return ""

# =====================================================
# TEMPLATE REMOVAL
# =====================================================
def remove_template(text, template):
    for line in template.splitlines():
        if len(line.strip()) > 10:
            text = text.replace(line, "")
    return text

# =====================================================
# CACHED EMBEDDINGS
# =====================================================
@st.cache_data(show_spinner=False)
def get_embedding(text):
    return model.encode(text, convert_to_tensor=True)

# =====================================================
# SIMILARITY
# =====================================================
def semantic_similarity(t1, t2):
    emb1 = get_embedding(t1)
    emb2 = get_embedding(t2)
    return max(util.cos_sim(emb1, emb2).item(), 0.0)

def overlapping_sentences(t1, t2):
    s1 = set(re.split(r'[.!?]', t1))
    s2 = set(re.split(r'[.!?]', t2))
    return [s.strip() for s in s1 & s2 if len(s.strip()) > 20]

# =====================================================
# UI
# =====================================================
st.subheader("📦 Upload Student Submissions (ZIP or Multiple Files)")

uploaded_files = st.file_uploader(
    "Accepted formats: ZIP, TXT, PDF, DOCX, PPTX",
    type=["zip", "txt", "pdf", "docx", "pptx"],
    accept_multiple_files=True
)

template_file = st.file_uploader(
    "Upload Template File (Optional)",
    type=["txt", "pdf", "docx"]
)

# =====================================================
# PROCESS
# =====================================================
if st.button("Run Plagiarism Check"):

    texts = {}
    template_text = ""

    if template_file:
        template_text = extract_file_text(template_file.read(), template_file.name)

    # Extract texts from uploaded files
    for f in uploaded_files:
        if f.name.lower().endswith(".zip"):
            with zipfile.ZipFile(f) as z:
                for name in z.namelist():
                    if name.lower().endswith((".txt", ".pdf", ".docx", ".pptx")):
                        with z.open(name) as extracted:
                            extracted_bytes = extracted.read()
                            texts[name] = extract_file_text(extracted_bytes, name)
        else:
            texts[f.name] = extract_file_text(f.read(), f.name)

    if len(texts) < 2:
        st.error("Please upload at least 2 submissions.")
        st.stop()

    # Remove template content
    cleaned_texts = {k: remove_template(v, template_text) for k, v in texts.items()}

    results = []
    pdf_elements = []
    styles = getSampleStyleSheet()

    # =====================================================
    # PROCESS PAIRS
    # =====================================================
    for (f1, t1), (f2, t2) in itertools.combinations(cleaned_texts.items(), 2):
        sim = semantic_similarity(t1, t2) * 100
        overlaps = overlapping_sentences(t1, t2)

        # Record similarity
        results.append({
            "File 1": f1,
            "File 2": f2,
            "Similarity (%)": round(sim, 2)
        })

        # PDF section
        pdf_elements.append(
            Paragraph(f"<b>{f1} vs {f2} — Similarity: {sim:.2f}%</b>", styles["Heading2"])
        )

        if overlaps:
            pdf_elements.append(Paragraph("<b>Overlapping Sentences:</b>", styles["Heading3"]))
            for s in overlaps:
                pdf_elements.append(Paragraph(f"<font color='red'>{s}</font>", styles["Normal"]))
        else:
            pdf_elements.append(Paragraph("No significant overlapping sentences detected.", styles["Normal"]))

        pdf_elements.append(Paragraph("<hr width='100%'/>", styles["Normal"]))
        pdf_elements.append(Paragraph(" ", styles["Normal"]))

    # =====================================================
    # SHOW RESULTS
    # =====================================================
    df = pd.DataFrame(results)
    st.success("Plagiarism analysis completed.")
    st.dataframe(df, use_container_width=True)

    # =====================================================
    # CSV DOWNLOAD IN-MEMORY
    # =====================================================
    csv_buffer = BytesIO()
    csv_buffer.write(df.to_csv(index=False).encode('utf-8'))
    csv_buffer.seek(0)

    st.download_button(
        "⬇ Download CSV Report (Lecturers)",
        data=csv_buffer,
        file_name="plagiarism_report.csv",
        mime="text/csv"
    )

    # =====================================================
    # PDF DOWNLOAD IN-MEMORY
    # =====================================================
    pdf_buffer = BytesIO()
    doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
    doc.build(pdf_elements)
    pdf_buffer.seek(0)

    st.download_button(
        "⬇ Download Highlighted PDF Report",
        data=pdf_buffer,
        file_name="plagiarism_report.pdf",
        mime="application/pdf"
    )
